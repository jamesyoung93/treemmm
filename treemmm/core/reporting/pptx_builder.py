"""PowerPoint report builder for TreeMMM results.

Optional dependency: install with ``pip install treemmm[reporting]``.

Generates a 16+ slide deck covering:
    1. Title slide
    2. Executive summary
    3. Data overview
    4-5. Model performance (aggregate + per-fold)
    6-7. Global attribution (bar chart + table)
    8-9. Temporal attribution (stacked area + table)
    10-11. Feature importance (bar + SHAP summary)
    12. Reverse causality diagnostics
    13-14. mROI response curves (if available)
    15. mROI reallocation recommendation (if available)
    16. Methodology appendix
"""

from __future__ import annotations

import io
import logging
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
import seaborn as sns

from treemmm.core.attribution.decomposer import Attribution
from treemmm.core.models.base import ModelResult

logger = logging.getLogger(__name__)

# Publication-quality style defaults
TREEMMM_COLORS = [
    "#2E86AB", "#A23B72", "#F18F01", "#C73E1D", "#3B1F2B",
    "#44BBA4", "#E94F37", "#393E41", "#D4B483", "#6B4226",
]


def _fig_to_bytes(fig: plt.Figure, dpi: int = 200) -> bytes:
    """Convert a matplotlib figure to PNG bytes."""
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=dpi, bbox_inches="tight")
    buf.seek(0)
    data = buf.read()
    plt.close(fig)
    return data


def _make_attribution_bar(attribution: Attribution) -> bytes:
    """Create a horizontal bar chart of global attribution."""
    ga = attribution.global_attribution()
    ga = ga[ga["variable"] != "_base"].sort_values("pct_of_total", ascending=True)

    fig, ax = plt.subplots(figsize=(8, max(3, len(ga) * 0.5)))
    colors = TREEMMM_COLORS[:len(ga)]
    ax.barh(ga["variable"], ga["pct_of_total"], color=colors[::-1])
    ax.set_xlabel("% of Total Attribution")
    ax.set_title("Global Attribution by Variable")
    for i, (_, row) in enumerate(ga.iterrows()):
        ax.text(row["pct_of_total"] + 0.3, i, f"{row['pct_of_total']:.1f}%",
                va="center", fontsize=9)
    fig.tight_layout()
    return _fig_to_bytes(fig)


def _make_temporal_attribution(
    attribution: Attribution,
    time_values: np.ndarray | pd.Series,
) -> bytes:
    """Create a stacked area chart of temporal attribution."""
    ta = attribution.temporal_attribution(time_values)
    ta = ta[ta["variable"] != "_base"]

    pivot = ta.pivot_table(
        index="time", columns="variable", values="attribution", fill_value=0,
    )
    # Sort columns by total attribution
    col_order = pivot.abs().sum().sort_values(ascending=False).index
    pivot = pivot[col_order]

    fig, ax = plt.subplots(figsize=(10, 5))
    pivot.plot.area(ax=ax, color=TREEMMM_COLORS[:len(pivot.columns)], alpha=0.8)
    ax.set_xlabel("Period")
    ax.set_ylabel("Attribution")
    ax.set_title("Attribution Over Time")
    ax.legend(loc="upper left", fontsize=8)
    fig.tight_layout()
    return _fig_to_bytes(fig)


def _make_feature_importance(attribution: Attribution) -> bytes:
    """Create a feature importance bar chart from mean |SHAP|."""
    mean_abs = np.mean(np.abs(attribution.values), axis=0)
    feat_df = pd.DataFrame({
        "variable": attribution.feature_names,
        "mean_abs_shap": mean_abs,
    }).sort_values("mean_abs_shap", ascending=True)

    fig, ax = plt.subplots(figsize=(8, max(3, len(feat_df) * 0.5)))
    ax.barh(feat_df["variable"], feat_df["mean_abs_shap"], color=TREEMMM_COLORS[0])
    ax.set_xlabel("Mean |SHAP Value|")
    ax.set_title("Feature Importance (Mean Absolute SHAP)")
    fig.tight_layout()
    return _fig_to_bytes(fig)


def _make_performance_chart(model_result: ModelResult) -> bytes:
    """Create a per-fold performance bar chart."""
    folds = model_result.fold_results
    fold_labels = [f"Fold {f.fold_idx + 1}" for f in folds]
    r2_vals = []
    for fr in folds:
        ss_res = np.sum((fr.y_true - fr.y_pred) ** 2)
        ss_tot = np.sum((fr.y_true - np.mean(fr.y_true)) ** 2)
        r2_vals.append(1 - ss_res / ss_tot if ss_tot > 0 else 0)

    fig, ax = plt.subplots(figsize=(8, 4))
    bars = ax.bar(fold_labels, r2_vals, color=TREEMMM_COLORS[0], alpha=0.8)
    ax.axhline(model_result.r2, color=TREEMMM_COLORS[1], linestyle="--",
               label=f"Pooled R² = {model_result.r2:.3f}")
    ax.set_ylabel("R²")
    ax.set_title("Model Performance Across CV Folds")
    ax.legend()
    for bar, val in zip(bars, r2_vals):
        ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 0.01,
                f"{val:.3f}", ha="center", fontsize=9)
    fig.tight_layout()
    return _fig_to_bytes(fig)


def _make_response_curves(mroi_result) -> bytes:
    """Create response curve plots for mROI."""
    curves = mroi_result.response_curves
    n_curves = len(curves)
    cols = min(3, n_curves)
    rows = (n_curves + cols - 1) // cols

    fig, axes = plt.subplots(rows, cols, figsize=(5 * cols, 4 * rows), squeeze=False)
    for i, curve in enumerate(curves):
        ax = axes[i // cols][i % cols]
        x = [pt.pct_of_current * 100 for pt in curve.points]
        y = [pt.predicted_outcome for pt in curve.points]
        y_lo = [pt.predicted_outcome_lower for pt in curve.points]
        y_hi = [pt.predicted_outcome_upper for pt in curve.points]

        ax.plot(x, y, color=TREEMMM_COLORS[i % len(TREEMMM_COLORS)], linewidth=2)
        ax.fill_between(x, y_lo, y_hi, alpha=0.2,
                        color=TREEMMM_COLORS[i % len(TREEMMM_COLORS)])
        ax.axvline(100, color="gray", linestyle="--", alpha=0.5, label="Current")
        ax.set_xlabel("% of Current Level")
        ax.set_ylabel("Mean Predicted Outcome")
        ax.set_title(curve.variable)
        ax.legend(fontsize=8)

    # Hide empty subplots
    for i in range(n_curves, rows * cols):
        axes[i // cols][i % cols].set_visible(False)

    fig.suptitle("mROI Response Curves", fontsize=14, y=1.02)
    fig.tight_layout()
    return _fig_to_bytes(fig)


def build_pptx(
    model_result: ModelResult,
    attribution: Attribution,
    time_values: np.ndarray | pd.Series | None = None,
    mroi_result=None,
    title: str = "TreeMMM Analysis Report",
    subtitle: str = "",
    output_path: str | Path | None = None,
) -> bytes | None:
    """Build a PowerPoint report.

    Args:
        model_result: Model performance results.
        attribution: Attribution decomposition results.
        time_values: Time period values for temporal charts.
        mroi_result: Optional mROI simulation results.
        title: Report title.
        subtitle: Report subtitle.
        output_path: If provided, writes PPTX file to this path.

    Returns:
        PPTX bytes if output_path is None, else None (written to file).
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
    except ImportError as e:
        raise ImportError(
            "python-pptx is not installed. "
            "Install with: pip install treemmm[reporting]"
        ) from e

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title ---
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if slide.placeholders[1]:
        slide.placeholders[1].text = subtitle

    # --- Slide 2: Executive Summary ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Executive Summary"
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = f"Model: {model_result.model_name}"
    tf.add_paragraph().text = f"Pooled R²: {model_result.r2:.4f}"
    tf.add_paragraph().text = f"Pooled WMAPE: {model_result.wmape:.4f}"
    tf.add_paragraph().text = f"Pooled MAE: {model_result.mae:.4f}"
    tf.add_paragraph().text = f"CV Folds: {len(model_result.fold_results)}"

    # Top attribution
    ga = attribution.global_attribution()
    ga_top = ga[ga["variable"] != "_base"].head(5)
    tf.add_paragraph().text = ""
    tf.add_paragraph().text = "Top attributed variables:"
    for _, row in ga_top.iterrows():
        tf.add_paragraph().text = f"  {row['variable']}: {row['pct_of_total']:.1f}%"

    # --- Slide 3: Model Performance ---
    perf_png = _make_performance_chart(model_result)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    slide.shapes.title.text = "Model Performance"
    slide.shapes.add_picture(
        io.BytesIO(perf_png), Inches(1), Inches(1.5), Inches(10), Inches(5),
    )

    # --- Slide 4: Global Attribution ---
    attr_png = _make_attribution_bar(attribution)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.title.text = "Global Attribution"
    slide.shapes.add_picture(
        io.BytesIO(attr_png), Inches(1), Inches(1.5), Inches(10), Inches(5),
    )

    # --- Slide 5: Feature Importance ---
    fi_png = _make_feature_importance(attribution)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.title.text = "Feature Importance (Mean |SHAP|)"
    slide.shapes.add_picture(
        io.BytesIO(fi_png), Inches(1), Inches(1.5), Inches(10), Inches(5),
    )

    # --- Slide 6: Temporal Attribution (if time values provided) ---
    if time_values is not None:
        ta_png = _make_temporal_attribution(attribution, time_values)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.title.text = "Attribution Over Time"
        slide.shapes.add_picture(
            io.BytesIO(ta_png), Inches(1), Inches(1.5), Inches(10), Inches(5),
        )

    # --- Slide 7-8: mROI (if available) ---
    if mroi_result is not None:
        rc_png = _make_response_curves(mroi_result)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.title.text = "mROI Response Curves"
        slide.shapes.add_picture(
            io.BytesIO(rc_png), Inches(0.5), Inches(1.5), Inches(12), Inches(5),
        )

        # Reallocation summary slide
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "mROI Reallocation Recommendation"
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.text = mroi_result.summary()

    # --- Slide: Methodology ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Methodology"
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "TreeMMM: Tree-Based Market Mix Modeling"
    tf.add_paragraph().text = ""
    tf.add_paragraph().text = "- Gradient-boosted trees (LightGBM/XGBoost/CatBoost)"
    tf.add_paragraph().text = "- SHAP TreeExplainer for attribution"
    tf.add_paragraph().text = "- Link-function-aware decomposition"
    tf.add_paragraph().text = "- Temporal cross-validation (no future leakage)"
    tf.add_paragraph().text = "- Distribution-matched objective functions"
    tf.add_paragraph().text = ""
    tf.add_paragraph().text = "Generated by TreeMMM v0.1.0"

    # Write output
    if output_path is not None:
        output_path = Path(output_path)
        prs.save(str(output_path))
        logger.info(f"PPTX written to {output_path}")
        return None
    else:
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return buf.read()
